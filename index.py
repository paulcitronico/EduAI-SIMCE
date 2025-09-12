import streamlit as st
import sqlite3
import hashlib
import os
from datetime import datetime, timedelta
import base64
from io import BytesIO
import subprocess
import platform
import tempfile
import secrets
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
# Importar librerías para previsualización de documentos
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

# Configuración inicial
st.set_page_config(page_title="Sistema de Login", layout="wide")

# Configuración de correo - Configura estos valores según tu servidor SMTP
SMTP_SERVER = "smtp.gmail.com"  # Cambiar según tu proveedor
SMTP_PORT = 587
SMTP_USERNAME = "jeanpachecotesista@gmail.com"  # Cambiar por tu correo
SMTP_PASSWORD = "ecea gcpe ygqw nyal"  # Cambiar por tu contraseña o contraseña de aplicación

# Crear directorios necesarios
if not os.path.exists("user_images"):
    os.makedirs("user_images")
if not os.path.exists("archivos_profesores"):
    os.makedirs("archivos_profesores")
if not os.path.exists("temp_conversion"):
    os.makedirs("temp_conversion")
if not os.path.exists("videos_tutoriales"):
    os.makedirs("videos_tutoriales")

# Funciones para la base de datos
def init_db():
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            nombre TEXT NOT NULL,
            apellido TEXT NOT NULL,
            email TEXT UNIQUE NOT NULL,
            password TEXT NOT NULL,
            imagen_path TEXT,
            rol TEXT DEFAULT 'alumno',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            reset_token TEXT,
            reset_token_expiration TEXT
        )
    ''')
    
    # Crear tabla de archivos con la columna categoria y ruta_pdf
    c.execute('''
        CREATE TABLE IF NOT EXISTS archivos (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            profesor_id INTEGER NOT NULL,
            nombre_archivo TEXT NOT NULL,
            ruta_archivo TEXT NOT NULL,
            tipo_archivo TEXT NOT NULL,
            categoria TEXT NOT NULL,
            ruta_pdf TEXT,
            fecha_subida TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (profesor_id) REFERENCES users(id)
        )
    ''')
    
    # Crear tabla para tutoriales
    c.execute('''
        CREATE TABLE IF NOT EXISTS tutoriales (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            profesor_id INTEGER NOT NULL,
            titulo TEXT NOT NULL,
            descripcion TEXT,
            url_youtube TEXT,
            ruta_video TEXT,
            fecha_creacion TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (profesor_id) REFERENCES users(id)
        )
    ''')
    
    # Verificar si la columna categoria existe, si no, agregarla
    c.execute("PRAGMA table_info(archivos)")
    columns = [column[1] for column in c.fetchall()]
    if 'categoria' not in columns:
        c.execute("ALTER TABLE archivos ADD COLUMN categoria TEXT")
    
    # Verificar si la columna ruta_pdf existe, si no, agregarla
    if 'ruta_pdf' not in columns:
        c.execute("ALTER TABLE archivos ADD COLUMN ruta_pdf TEXT")
    
    # Verificar si las columnas de reseteo existen
    c.execute("PRAGMA table_info(users)")
    user_columns = [column[1] for column in c.fetchall()]
    if 'reset_token' not in user_columns:
        c.execute("ALTER TABLE users ADD COLUMN reset_token TEXT")
    if 'reset_token_expiration' not in user_columns:
        c.execute("ALTER TABLE users ADD COLUMN reset_token_expiration TEXT")
    
    # Crear usuario administrador si no existe
    c.execute("SELECT * FROM users WHERE username = 'admin'")
    if not c.fetchone():
        hashed_pw = hash_password('admin123')
        c.execute('''
            INSERT INTO users (username, nombre, apellido, email, password, rol)
            VALUES (?, ?, ?, ?, ?, ?)
        ''', ('admin', 'Administrador', 'Sistema', 'admin@example.com', hashed_pw, 'admin'))
    
    conn.commit()
    conn.close()

def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()

def send_email(to_email, subject, body):
    try:
        msg = MIMEMultipart()
        msg['From'] = SMTP_USERNAME
        msg['To'] = to_email
        msg['Subject'] = subject
        
        msg.attach(MIMEText(body, 'plain'))
        
        server = smtplib.SMTP(SMTP_SERVER, SMTP_PORT)
        server.starttls()
        server.login(SMTP_USERNAME, SMTP_PASSWORD)
        text = msg.as_string()
        server.sendmail(SMTP_USERNAME, to_email, text)
        server.quit()
        return True
    except Exception as e:
        st.error(f"Error al enviar el correo: {str(e)}")
        return False

def show_all_emails():
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute("SELECT email FROM users")
    emails = c.fetchall()
    conn.close()
    return [email[0] for email in emails]

def add_user(username, nombre, apellido, email, password, imagen_path=None):
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    try:
        hashed_pw = hash_password(password)
        c.execute('''
            INSERT INTO users (username, nombre, apellido, email, password, imagen_path, rol)
            VALUES (?, ?, ?, ?, ?, ?, 'alumno')
        ''', (username, nombre, apellido, email, hashed_pw, imagen_path))
        conn.commit()
        
        # Enviar correo de bienvenida
        subject = "Bienvenido al Sistema"
        body = f"""
        Hola {nombre},
        
        Te damos la bienvenida a nuestro sistema educativo.
        
        Tu cuenta ha sido creada exitosamente con los siguientes datos:
        - Usuario: {username}
        - Nombre: {nombre} {apellido}
        - Correo: {email}
        
        Para iniciar sesión, visita nuestra aplicación e ingresa tus credenciales.
        
        Si tienes alguna pregunta, no dudes en contactarnos.
        
        Saludos,
        El equipo del sistema
        """
        
        if send_email(email, subject, body):
            st.success("¡Cuenta creada exitosamente! Revisa tu correo para el mensaje de bienvenida.")
        else:
            st.warning("Tu cuenta fue creada, pero hubo un error al enviar el correo de bienvenida.")
        
        return True
    except sqlite3.IntegrityError:
        return False
    finally:
        conn.close()

def verify_user(username, password):
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    hashed_pw = hash_password(password)
    c.execute('SELECT * FROM users WHERE username = ? AND password = ?', (username, hashed_pw))
    user = c.fetchone()
    conn.close()
    return user

def update_user_role(user_id, new_role):
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute("UPDATE users SET rol = ? WHERE id = ?", (new_role, user_id))
    conn.commit()
    conn.close()

def get_all_users():
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute("SELECT id, username, nombre, apellido, email, rol FROM users")
    users = c.fetchall()
    conn.close()
    return users

def delete_user(user_id):
    """Elimina un usuario de la base de datos y su imagen si existe"""
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    
    # Obtener información del usuario
    c.execute("SELECT imagen_path FROM users WHERE id = ?", (user_id,))
    result = c.fetchone()
    
    if result:
        imagen_path = result[0]
        
        # Eliminar imagen si existe
        try:
            if imagen_path and os.path.exists(imagen_path):
                os.remove(imagen_path)
        except Exception as e:
            conn.close()
            return False, f"Error al eliminar la imagen: {str(e)}"
        
        # Eliminar usuario de la base de datos
        c.execute("DELETE FROM users WHERE id = ?", (user_id,))
        conn.commit()
        conn.close()
        return True, "Usuario eliminado exitosamente"
    else:
        conn.close()
        return False, "Usuario no encontrado en la base de datos"

def generate_reset_token(email):
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    
    # Verificar si el email existe (insensible a mayúsculas/minúsculas)
    c.execute("SELECT id, nombre FROM users WHERE LOWER(email) = LOWER(?)", (email,))
    user = c.fetchone()
    
    if not user:
        conn.close()
        return False, "El correo electrónico no está registrado en el sistema."
    
    user_id, nombre = user
    
    # Generar token único
    token = secrets.token_urlsafe(32)
    expiration = datetime.now() + timedelta(hours=1)  # Token válido por 1 hora
    
    # Guardar token en la base de datos (convertir a string ISO)
    expiration_str = expiration.isoformat()
    c.execute("UPDATE users SET reset_token = ?, reset_token_expiration = ? WHERE id = ?", 
              (token, expiration_str, user_id))
    conn.commit()
    conn.close()
    
    # Enviar correo con el token
    subject = "Restablecimiento de Contraseña"
    body = f"""
    Hola {nombre},
    
    Hemos recibido una solicitud para restablecer tu contraseña.
    
    Para continuar, haz clic en el siguiente enlace o copia y pégalo en tu navegador:
    http://localhost:8501/?reset_token={token}
    
    Este enlace expirará en 1 hora.
    
    Si no solicitaste restablecer tu contraseña, ignora este mensaje.
    
    Saludos,
    El equipo del sistema
    """
    
    # Intentar enviar el correo y capturar errores
    try:
        if send_email(email, subject, body):
            return True, "Se ha enviado un enlace de restablecimiento a tu correo electrónico."
        else:
            return False, "Hubo un error al enviar el correo. Inténtalo de nuevo."
    except Exception as e:
        return False, f"Error al enviar el correo: {str(e)}"

def reset_password_with_token(token, new_password):
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    
    # Verificar token
    c.execute("SELECT id, reset_token_expiration FROM users WHERE reset_token = ?", (token,))
    user = c.fetchone()
    
    if not user:
        conn.close()
        return False, "Token inválido o expirado."
    
    user_id, expiration_str = user
    
    try:
        # Convertir string a datetime
        expiration = datetime.fromisoformat(expiration_str)
    except ValueError:
        conn.close()
        return False, "Formato de token inválido."
    
    if datetime.now() > expiration:
        conn.close()
        return False, "Token expirado. Solicita un nuevo restablecimiento de contraseña."
    
    # Actualizar contraseña
    hashed_pw = hash_password(new_password)
    c.execute("UPDATE users SET password = ?, reset_token = NULL, reset_token_expiration = NULL WHERE id = ?", 
              (hashed_pw, user_id))
    conn.commit()
    conn.close()
    
    return True, "Contraseña restablecida exitosamente."

def eliminar_archivo(archivo_id):
    """Elimina un archivo de la base de datos y del sistema de archivos"""
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    
    # Obtener información del archivo
    c.execute("SELECT ruta_archivo, ruta_pdf FROM archivos WHERE id = ?", (archivo_id,))
    resultado = c.fetchone()
    
    if resultado:
        ruta_archivo, ruta_pdf = resultado
        
        # Eliminar archivos físicos
        try:
            if os.path.exists(ruta_archivo):
                os.remove(ruta_archivo)
            
            if ruta_pdf and os.path.exists(ruta_pdf):
                os.remove(ruta_pdf)
        except Exception as e:
            conn.close()
            return False, f"Error al eliminar archivos físicos: {str(e)}"
        
        # Eliminar registro de la base de datos
        c.execute("DELETE FROM archivos WHERE id = ?", (archivo_id,))
        conn.commit()
        conn.close()
        return True, "Archivo eliminado exitosamente"
    else:
        conn.close()
        return False, "Archivo no encontrado en la base de datos"

def convertir_pptx_a_pdf(ruta_entrada, ruta_salida):
    """
    Convierte un archivo PPTX a PDF usando LibreOffice o Microsoft Office
    """
    try:
        # Determinar el sistema operativo
        sistema = platform.system()
        
        if sistema == 'Windows':
            # En Windows, intentar usar Microsoft Office si está disponible
            try:
                import win32com.client
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                deck = powerpoint.Presentations.Open(ruta_entrada)
                deck.SaveAs(ruta_salida, 32)  # 32 es el formato PDF
                deck.Close()
                powerpoint.Quit()
                return True
            except:
                # Si falla Microsoft Office, intentar con LibreOffice
                pass
        
        # Usar LibreOffice (funciona en Windows, Linux y macOS)
        # Crear un directorio temporal para la conversión
        temp_dir = tempfile.mkdtemp()
        temp_salida = os.path.join(temp_dir, "temp.pdf")
        
        # Comando para convertir usando LibreOffice
        comando = [
            'libreoffice', '--headless', '--convert-to', 'pdf', 
            '--outdir', temp_dir, ruta_entrada
        ]
        
        # Ejecutar el comando
        subprocess.run(comando, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        
        # Mover el archivo convertido a la ruta de salida
        if os.path.exists(temp_salida):
            os.rename(temp_salida, ruta_salida)
            return True
        
        return False
    except Exception as e:
        st.error(f"Error al convertir PPTX a PDF: {str(e)}")
        return False

def guardar_archivo(profesor_id, archivo, categoria):
    # Guardar archivo en disco
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    nombre_archivo = f"{timestamp}_{archivo.name}"
    ruta_archivo = f"archivos_profesores/{nombre_archivo}"
    
    with open(ruta_archivo, "wb") as f:
        f.write(archivo.getbuffer())
    
    ruta_pdf = None
    
    # Si es un archivo PPTX, intentar convertirlo a PDF
    if archivo.type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                       'application/vnd.ms-powerpoint',
                       'application/vnd.ms-powerpoint.presentation.macroEnabled.12']:
        try:
            # Ruta para el PDF convertido
            nombre_pdf = f"{timestamp}_{os.path.splitext(archivo.name)[0]}.pdf"
            ruta_pdf = f"archivos_profesores/{nombre_pdf}"
            
            # Intentar convertir a PDF
            if convertir_pptx_a_pdf(ruta_archivo, ruta_pdf):
                st.success(f"Archivo convertido a PDF: {nombre_pdf}")
            else:
                ruta_pdf = None
                st.warning("No se pudo convertir el archivo a PDF. Se mostrará el formato original.")
        except Exception as e:
            ruta_pdf = None
            st.warning(f"Error al convertir el archivo a PDF: {str(e)}. Se mostrará el formato original.")
    
    # Guardar información en la base de datos
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute('''
        INSERT INTO archivos (profesor_id, nombre_archivo, ruta_archivo, tipo_archivo, categoria, ruta_pdf)
        VALUES (?, ?, ?, ?, ?, ?)
    ''', (profesor_id, archivo.name, ruta_archivo, archivo.type, categoria, ruta_pdf))
    conn.commit()
    conn.close()
    
    return ruta_archivo

def obtener_archivos(categoria=None, profesor_id=None):
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    
    query = '''
        SELECT a.id, a.nombre_archivo, a.ruta_archivo, a.tipo_archivo, a.fecha_subida, u.nombre, u.apellido, a.categoria, a.ruta_pdf
        FROM archivos a
        JOIN users u ON a.profesor_id = u.id
    '''
    
    params = []
    conditions = []
    
    if categoria:
        conditions.append("a.categoria = ?")
        params.append(categoria)
    
    if profesor_id:
        conditions.append("a.profesor_id = ?")
        params.append(profesor_id)
    
    if conditions:
        query += " WHERE " + " AND ".join(conditions)
    
    query += " ORDER BY a.fecha_subida DESC"
    
    c.execute(query, params)
    archivos = c.fetchall()
    conn.close()
    return archivos

def agregar_tutorial(profesor_id, titulo, descripcion, url_youtube=None, ruta_video=None):
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    try:
        c.execute('''
            INSERT INTO tutoriales (profesor_id, titulo, descripcion, url_youtube, ruta_video)
            VALUES (?, ?, ?, ?, ?)
        ''', (profesor_id, titulo, descripcion, url_youtube, ruta_video))
        conn.commit()
        return True
    except Exception as e:
        st.error(f"Error al agregar tutorial: {str(e)}")
        return False
    finally:
        conn.close()

def obtener_tutoriales():
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute('''
        SELECT t.id, t.titulo, t.descripcion, t.url_youtube, t.ruta_video, t.fecha_creacion, u.nombre, u.apellido
        FROM tutoriales t
        JOIN users u ON t.profesor_id = u.id
        ORDER BY t.fecha_creacion DESC
    ''')
    tutoriales = c.fetchall()
    conn.close()
    return tutoriales

def eliminar_tutorial(tutorial_id):
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    
    # Obtener información del tutorial
    c.execute("SELECT ruta_video FROM tutoriales WHERE id = ?", (tutorial_id,))
    resultado = c.fetchone()
    
    if resultado:
        ruta_video = resultado[0]
        
        # Eliminar archivo de video si existe
        try:
            if ruta_video and os.path.exists(ruta_video):
                os.remove(ruta_video)
        except Exception as e:
            conn.close()
            return False, f"Error al eliminar archivo de video: {str(e)}"
        
        # Eliminar registro de la base de datos
        c.execute("DELETE FROM tutoriales WHERE id = ?", (tutorial_id,))
        conn.commit()
        conn.close()
        return True, "Tutorial eliminado exitosamente"
    else:
        conn.close()
        return False, "Tutorial no encontrado en la base de datos"

def mostrar_pdf(ruta_archivo):
    with open(ruta_archivo, "rb") as f:
        base64_pdf = base64.b64encode(f.read()).decode('utf-8')
    
    pdf_display = f'<iframe src="data:application/pdf;base64,{base64_pdf}" width="100%" height="800" type="application/pdf"></iframe>'
    st.markdown(pdf_display, unsafe_allow_html=True)

def mostrar_docx_completo(ruta_archivo):
    if not DOCX_AVAILABLE:
        st.error("La librería 'python-docx' no está instalada. No se puede previsualizar el archivo DOCX.")
        return
    
    try:
        doc = Document(ruta_archivo)
        
        # Procesar cada párrafo
        for para in doc.paragraphs:
            if para.text.strip():
                st.write(para.text)
        
        # Procesar tablas
        for table in doc.tables:
            st.write("---")
            for i, row in enumerate(table.rows):
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text)
                st.write(" | ".join(row_data))
            st.write("---")
        
        # Procesar imágenes
        image_count = 0
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_path = os.path.join(os.path.dirname(ruta_archivo), rel.target_ref)
                    if os.path.exists(image_path):
                        st.write("---")
                        st.write(f"Imagen {image_count + 1}:")
                        st.image(image_path, width=600)
                        image_count += 1
                except:
                    pass
                
    except Exception as e:
        st.error(f"No se pudo previsualizar el archivo DOCX: {str(e)}")

def mostrar_pptx_completo(ruta_archivo):
    if not PPTX_AVAILABLE:
        st.error("La librería 'python-pptx' no está instalada. No se puede previsualizar el archivo PPTX.")
        return
    
    try:
        prs = Presentation(ruta_archivo)
        
        # Procesar cada diapositiva
        for i, slide in enumerate(prs.slides):
            # Crear un contenedor para la diapositiva
            with st.container():
                st.markdown(f"### Diapositiva {i+1}")
                
                # Extraer y mostrar el título de la diapositiva si existe
                slide_title = ""
                for shape in slide.shapes:
                    # Verificar si la forma tiene placeholder_format antes de acceder a ella
                    if hasattr(shape, "placeholder_format"):
                        if shape.placeholder_format.type == 1:  # 1 es el título
                            slide_title = shape.text if shape.text else f"Diapositiva {i+1}"
                            break
                
                if slide_title:
                    st.markdown(f"#### {slide_title}")
                
                # Contenedor para el contenido de la diapositiva
                slide_content = []
                images_in_slide = []
                
                # Procesar cada forma en la diapositiva
                for shape in slide.shapes:
                    # Si es un cuadro de texto
                    if hasattr(shape, "text") and shape.text:
                        slide_content.append(shape.text)
                    
                    # Si es una imagen
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            
                            # Guardar la imagen temporalmente
                            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
                            temp_img_path = f"temp_conversion/{timestamp}_slide{i}_img.png"
                            with open(temp_img_path, "wb") as img_file:
                                img_file.write(image_bytes)
                            
                            images_in_slide.append(temp_img_path)
                        except Exception as e:
                            st.error(f"Error al procesar imagen: {str(e)}")
                    
                    # Si es un grupo de formas
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, "text") and sub_shape.text:
                                slide_content.append(sub_shape.text)
                            elif sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                try:
                                    image = sub_shape.image
                                    image_bytes = image.blob
                                    
                                    # Guardar la imagen temporalmente
                                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
                                    temp_img_path = f"temp_conversion/{timestamp}_slide{i}_img.png"
                                    with open(temp_img_path, "wb") as img_file:
                                        img_file.write(image_bytes)
                                    
                                    images_in_slide.append(temp_img_path)
                                except Exception as e:
                                    st.error(f"Error al procesar imagen en grupo: {str(e)}")
                
                # Mostrar el contenido de texto
                if slide_content:
                    for text in slide_content:
                        if text.strip():  # Solo mostrar si no está vacío
                            st.write(text)
                
                # Mostrar las imágenes
                if images_in_slide:
                    st.write("---")
                    st.write("Imágenes en la diapositiva:")
                    
                    # Determinar el número de columnas basado en la cantidad de imágenes
                    num_cols = min(len(images_in_slide), 3)
                    cols = st.columns(num_cols)
                    
                    for j, img_path in enumerate(images_in_slide):
                        with cols[j % num_cols]:
                            try:
                                st.image(img_path, width=300)
                            except Exception as e:
                                st.error(f"No se pudo mostrar la imagen: {str(e)}")
                
                # Separador entre diapositivas
                st.markdown("---")
                
    except Exception as e:
        st.error(f"No se pudo previsualizar el archivo PPTX: {str(e)}")

def mostrar_imagen(ruta_archivo):
    if not PIL_AVAILABLE:
        st.error("La librería 'PIL' no está instalada. No se puede previsualizar la imagen.")
        return
    
    try:
        image = Image.open(ruta_archivo)
        st.image(image, width=800)
    except Exception as e:
        st.error(f"No se pudo previsualizar la imagen: {str(e)}")

def mostrar_vista_previa(ruta_archivo, tipo_archivo, ruta_pdf=None):
    if tipo_archivo == 'application/pdf':
        st.write("Vista previa del PDF:")
        mostrar_pdf(ruta_archivo)
    elif tipo_archivo in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 
                         'application/msword',
                         'application/vnd.ms-word.document.macroEnabled.12']:
        st.write("Vista previa del documento:")
        mostrar_docx_completo(ruta_archivo)
    elif tipo_archivo in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         'application/vnd.ms-powerpoint',
                         'application/vnd.ms-powerpoint.presentation.macroEnabled.12']:
        # Si existe un PDF convertido, mostrarlo
        if ruta_pdf and os.path.exists(ruta_pdf):
            st.write("Vista previa de la presentación (convertida a PDF):")
            mostrar_pdf(ruta_pdf)
        else:
            st.write("Vista previa de la presentación:")
            mostrar_pptx_completo(ruta_archivo)
    elif tipo_archivo.startswith('image/'):
        st.write("Vista previa de la imagen:")
        mostrar_imagen(ruta_archivo)
    else:
        st.write("Vista previa no disponible para este tipo de archivo")

def extraer_id_youtube(url):
    """Extrae el ID del video de YouTube de una URL"""
    import re
    patrones = [
        r'(?:youtube\.com\/watch\?v=|\/v\/|youtu\.be\/)([^&]+)',
        r'(?:youtube\.com\/embed\/)([^&]+)',
        r'(?:youtube\.com\/v\/)([^&]+)'
    ]
    
    for patron in patrones:
        match = re.search(patron, url)
        if match:
            return match.group(1)
    return None

# Inicializar base de datos
init_db()

# Estado de sesión
if 'logged_in' not in st.session_state:
    st.session_state.logged_in = False
if 'current_user' not in st.session_state:
    st.session_state.current_user = None
if 'page' not in st.session_state:
    st.session_state.page = "login"
if 'admin_subpage' not in st.session_state:
    st.session_state.admin_subpage = "dashboard"
if 'selected_category' not in st.session_state:
    st.session_state.selected_category = "Unidades"
if 'reset_token' not in st.session_state:
    st.session_state.reset_token = None
if 'password_reset_success' not in st.session_state:
    st.session_state.password_reset_success = False

# Verificar si hay un token de restablecimiento en la URL
query_params = st.query_params
if 'reset_token' in query_params and not st.session_state.logged_in and not st.session_state.password_reset_success:
    st.session_state.reset_token = query_params['reset_token']
    st.session_state.page = "reset_password_confirm"

def go_to_login():
    st.query_params.clear()
    st.session_state.page = "login"
    st.session_state.reset_token = None
    st.rerun()

# Interfaz principal
def login_page():
    st.title("Iniciar Sesión")
    with st.form("login_form"):
        username = st.text_input("Username")
        password = st.text_input("Password", type="password")
        submitted = st.form_submit_button("Iniciar Sesión")
        
        if submitted:
            user = verify_user(username, password)
            if user:
                st.session_state.logged_in = True
                st.session_state.current_user = user
                st.success(f"Bienvenido {user[2]}!")
                st.rerun()
            else:
                st.error("Usuario o contraseña incorrectos")
    
    st.markdown("---")
    col1, col2 = st.columns(2)
    with col1:
        if st.button("¿No tienes cuenta? Regístrate"):
            st.session_state.page = "register"
            st.rerun()
    with col2:
        if st.button("¿Olvidaste tu contraseña?"):
            st.session_state.page = "reset_password"
            st.rerun()

def register_page():
    st.title("Crear Nueva Cuenta")
    with st.form("register_form"):
        username = st.text_input("Nombre de usuario")
        nombre = st.text_input("Nombre")
        apellido = st.text_input("Apellido")
        email = st.text_input("E-mail")
        password = st.text_input("Password", type="password")
        confirm_password = st.text_input("Confirmar Password", type="password")
        imagen = st.file_uploader("Subir una imagen", type=["jpg", "png", "jpeg"])
        
        submitted = st.form_submit_button("Registrarse")
        
        if submitted:
            # Validaciones
            if not all([username, nombre, apellido, email, password]):
                st.error("Por favor completa todos los campos")
            elif password != confirm_password:
                st.error("Las contraseñas no coinciden")
            else:
                # Guardar imagen si se subió
                imagen_path = None
                if imagen:
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    imagen_path = f"user_images/{username}_{timestamp}.jpg"
                    with open(imagen_path, "wb") as f:
                        f.write(imagen.getbuffer())
                
                # Agregar usuario a la base de datos (rol 'alumno' por defecto)
                if add_user(username, nombre, apellido, email, password, imagen_path):
                    go_to_login()
                else:
                    st.error("El nombre de usuario o email ya están en uso")
    
    st.markdown("---")
    if st.button("¿Ya tienes cuenta? Inicia sesión"):
        go_to_login()

def reset_password_page():
    st.title("Restablecer Contraseña")
    st.write("Ingresa tu correo electrónico y te enviaremos un enlace para restablecer tu contraseña.")
    
    # Botón para mostrar todos los correos (solo para depuración)
    if st.button("Mostrar todos los correos (depuración)"):
        all_emails = show_all_emails()
        st.write("Correos registrados:")
        for email in all_emails:
            st.write(f"- {email}")
    
    with st.form("reset_password_form"):
        email = st.text_input("Correo electrónico")
        submitted = st.form_submit_button("Enviar Enlace")
        
        if submitted:
            # Verificar si el correo existe (para depuración)
            conn = sqlite3.connect('users.db')
            c = conn.cursor()
            c.execute("SELECT email FROM users WHERE LOWER(email) = LOWER(?)", (email,))
            user_email = c.fetchone()
            conn.close()
            
            if user_email:
                st.info(f"Correo encontrado: {user_email[0]}")
            else:
                st.warning("Correo no encontrado en la base de datos.")
            
            success, message = generate_reset_token(email)
            if success:
                st.success(message)
                go_to_login()
            else:
                st.error(message)
    st.markdown("---")
    if st.button("¿Ya tienes cuenta? Inicia sesión"):
        go_to_login()

    

def reset_password_confirm_page():
    st.title("Restablecer Contraseña")
    token = st.session_state.reset_token
    
    # Si la contraseña ya se restableció con éxito, mostrar mensaje y botón para ir al login
    if st.session_state.password_reset_success:
        st.success("¡Contraseña restablecida exitosamente!")
        st.write("Ya puedes iniciar sesión con tu nueva contraseña.")
        
        if st.button("¿Ya tienes cuenta? Inicia sesión"):
            # Limpiar todos los estados relacionados con el restablecimiento
            st.session_state.password_reset_success = False
            go_to_login()
        return
    
    with st.form("reset_password_confirm_form"):
        new_password = st.text_input("Nueva Contraseña", type="password")
        confirm_password = st.text_input("Confirmar Nueva Contraseña", type="password")
        submitted = st.form_submit_button("Restablecer Contraseña")
        
        if submitted:
            if new_password != confirm_password:
                st.error("Las contraseñas no coinciden.")
            else:
                success, message = reset_password_with_token(token, new_password)
                if success:
                    # Marcar que la contraseña se restableció con éxito
                    st.session_state.password_reset_success = True
                    st.rerun()  # Recargar para mostrar el mensaje de éxito
                else:
                    st.error(message)
    
    st.markdown("---")
    if st.button("¿Ya tienes cuenta? Inicia sesión"):
        go_to_login()

def admin_dashboard():
    st.title("Panel de Administración - Dashboard")
    user = st.session_state.current_user
    st.write(f"Bienvenido, {user[2]} {user[3]}")
    st.write(f"Email: {user[4]}")
    st.write(f"Rol: {user[7]}")
    
    if user[6]:
        st.image(user[6], width=200)
    
    st.subheader("Opciones de Administración")
    if st.button("Gestionar Usuarios"):
        st.session_state.admin_subpage = "user_management"
        st.rerun()
    
    if st.button("Cerrar Sesión"):
        st.session_state.logged_in = False
        st.session_state.current_user = None
        st.rerun()

def admin_user_management():
    st.title("Gestión de Usuarios")
    
    # Mostrar todos los usuarios
    users = get_all_users()
    st.subheader("Lista de Usuarios")
    
    for user in users:
        user_id, username, nombre, apellido, email, rol = user
        
        # No mostrar al administrador en la lista
        if username == 'admin':
            continue
            
        col1, col2, col3, col4, col5, col6 = st.columns([1, 2, 2, 2, 2, 1])
        with col1:
            st.write(user_id)
        with col2:
            st.write(username)
        with col3:
            st.write(f"{nombre} {apellido}")
        with col4:
            st.write(email)
        with col5:
            # Selector de rol solo para el administrador
            new_rol = st.selectbox(
                "Rol",
                ['alumno', 'profesor'],
                index=0 if rol == 'alumno' else 1,
                key=f"rol_{user_id}"
            )
            
            if new_rol != rol:
                update_user_role(user_id, new_rol)
                st.success(f"Rol de {username} actualizado a {new_rol}")
                st.rerun()
        with col6:
            # Botón para eliminar usuario
            if st.button("🗑️", key=f"delete_{user_id}"):
                success, message = delete_user(user_id)
                if success:
                    st.success(message)
                    st.rerun()
                else:
                    st.error(message)
    
    if st.button("Volver al Dashboard"):
        st.session_state.admin_subpage = "dashboard"
        st.rerun()

def profesor_sidebar():
    # Barra lateral simplificada
    user = st.session_state.current_user
    st.sidebar.image(user[6] if user[6] else "https://via.placeholder.com/150", width=150)
    st.sidebar.write(f"**{user[2]} {user[3]}**")
    st.sidebar.write(f"Rol: {user[7]}")
    
    st.sidebar.markdown("---")
    
    if st.sidebar.button("Cerrar Sesión"):
        st.session_state.logged_in = False
        st.session_state.current_user = None
        st.rerun()

def profesor_page():
    user = st.session_state.current_user
    
    # Mostrar barra lateral simplificada
    profesor_sidebar()
    
    # Contenido principal
    st.title(f"Bienvenido, {user[2]} {user[3]}")
    
    # Barra de navegación en el contenido principal
    st.subheader("Navegación")
    col1, col2, col3, col4, col5 = st.columns(5)
    with col1:
        if st.button("Unidades", key="main_unidades"):
            st.session_state.selected_category = "Unidades"
            st.rerun()
    with col2:
        if st.button("Cuestionarios IA", key="main_cuestionarios"):
            st.session_state.selected_category = "Cuestionarios IA"
            st.rerun()
    with col3:
        if st.button("Revisiones", key="main_revisiones"):
            st.session_state.selected_category = "Revisiones"
            st.rerun()
    with col4:
        if st.button("Tutoriales", key="main_tutoriales"):
            st.session_state.selected_category = "Tutoriales"
            st.rerun()
    with col5:
        if st.button("Guías", key="main_guias"):
            st.session_state.selected_category = "Guías"
            st.rerun()
    
    # Contenido según la categoría seleccionada
    category = st.session_state.selected_category
    
    if category == "Cuestionarios IA":
        st.title(f"{category} - En Construcción")
        st.write("Esta sección está en desarrollo. Pronto estará disponible.")
    
    elif category == "Tutoriales":
        st.title("Gestión de Tutoriales")
        
        # Formulario para agregar tutoriales
        with st.form("form_tutorial"):
            st.subheader("Agregar Nuevo Tutorial")
            titulo = st.text_input("Título del tutorial")
            descripcion = st.text_area("Descripción")
            
            col1, col2 = st.columns(2)
            with col1:
                url_youtube = st.text_input("URL de YouTube (opcional)", placeholder="https://www.youtube.com/watch?v=...")
            with col2:
                video_file = st.file_uploader("O subir video MP4 (opcional)", type=["mp4"])
            
            submitted = st.form_submit_button("Agregar Tutorial")
            
            if submitted:
                if not titulo:
                    st.error("El título es obligatorio")
                elif not url_youtube and not video_file:
                    st.error("Debes proporcionar una URL de YouTube o subir un video MP4")
                else:
                    ruta_video = None
                    if video_file:
                        # Guardar video
                        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                        nombre_video = f"{timestamp}_{video_file.name}"
                        ruta_video = f"videos_tutoriales/{nombre_video}"
                        with open(ruta_video, "wb") as f:
                            f.write(video_file.getbuffer())
                    
                    if agregar_tutorial(user[0], titulo, descripcion, url_youtube, ruta_video):
                        st.success("Tutorial agregado exitosamente!")
                        st.rerun()
                    else:
                        st.error("Error al agregar el tutorial")
        
        # Mostrar tutoriales existentes
        st.subheader("Tutoriales Existentes")
        tutoriales = obtener_tutoriales()
        
        if not tutoriales:
            st.info("No hay tutoriales disponibles.")
        else:
            for tut in tutoriales:
                tut_id, titulo, descripcion, url_youtube, ruta_video, fecha_creacion, nombre_profesor, apellido_profesor = tut
                
                with st.expander(f"{titulo} - {fecha_creacion}"):
                    st.write(f"**Descripción:** {descripcion}")
                    st.write(f"**Creado por:** {nombre_profesor} {apellido_profesor}")
                    
                    # Mostrar video según el tipo
                    if url_youtube:
                        st.write("**Video de YouTube:**")
                        video_id = extraer_id_youtube(url_youtube)
                        if video_id:
                            st.video(f"https://www.youtube.com/watch?v={video_id}")
                        else:
                            st.error("URL de YouTube no válida")
                    
                    elif ruta_video and os.path.exists(ruta_video):
                        st.write("**Video subido:**")
                        st.video(ruta_video)
                    
                    # Botones de acción
                    col1, col2 = st.columns(2)
                    with col1:
                        if url_youtube:
                            st.link_button("🔗 Abrir en YouTube", url_youtube)
                        elif ruta_video:
                            with open(ruta_video, "rb") as file:
                                st.download_button(
                                    label="📥 Descargar Video",
                                    data=file,
                                    file_name=os.path.basename(ruta_video),
                                    mime="video/mp4"
                                )
                    with col2:
                        if st.button("🗑️ Eliminar", key=f"eliminar_tut_{tut_id}"):
                            success, message = eliminar_tutorial(tut_id)
                            if success:
                                st.success(message)
                                st.rerun()
                            else:
                                st.error(message)
    
    elif category == "Revisiones":
        # SECCIÓN MODIFICADA: Solo mostrar botón de descarga para Revisiones
        st.title("Revisiones - Descarga de Material")
        
        # Mostrar archivos existentes en Revisiones
        st.subheader("Material de Revisiones Disponible")
        archivos = obtener_archivos(categoria=category, profesor_id=user[0])
        
        if not archivos:
            st.info("No hay material de revisiones disponible actualmente.")
        else:
            for arch in archivos:
                arch_id, nombre_archivo, ruta_archivo, tipo_archivo, fecha_subida, nombre_profesor, apellido_profesor, categoria_archivo, ruta_pdf = arch
                
                # Crear expander para cada archivo
                with st.expander(f"{nombre_archivo} - {fecha_subida}"):
                    st.write(f"Tipo: {tipo_archivo}")
                    st.write(f"Categoría: {categoria_archivo}")
                    
                    # Vista previa según tipo de archivo
                    mostrar_vista_previa(ruta_archivo, tipo_archivo, ruta_pdf)
                    
                    # Solo botón de descarga (sin opción de eliminar)
                    with open(ruta_archivo, "rb") as file:
                        st.download_button(
                            label="📥 Descargar Archivo",
                            data=file,
                            file_name=nombre_archivo,
                            mime=tipo_archivo,
                            use_container_width=True
                        )
    else:
        # Para otras categorías (Unidades y Guías) mantener la funcionalidad original
        st.title(f"Gestión de {category}")
        
        # Subir archivos para Unidades y Guías
        archivo = st.file_uploader(f"Subir nuevo archivo a {category}", type=["pdf", "docx", "pptx", "jpg", "png", "jpeg"])
        
        if archivo is not None:
            if st.button("Subir Archivo"):
                ruta = guardar_archivo(user[0], archivo, category)
                st.success(f"Archivo '{archivo.name}' subido exitosamente a {category}")
                st.rerun()
        
        # Mostrar archivos existentes de la categoría seleccionada
        st.subheader(f"Archivos en {category}")
        archivos = obtener_archivos(categoria=category, profesor_id=user[0])
        
        if not archivos:
            st.write(f"No has subido ningún archivo a {category} aún.")
        else:
            for arch in archivos:
                arch_id, nombre_archivo, ruta_archivo, tipo_archivo, fecha_subida, nombre_profesor, apellido_profesor, categoria_archivo, ruta_pdf = arch
                
                # Crear expander para cada archivo
                with st.expander(f"{nombre_archivo} - {fecha_subida}"):
                    st.write(f"Tipo: {tipo_archivo}")
                    st.write(f"Categoría: {categoria_archivo}")
                    
                    # Vista previa según tipo de archivo
                    mostrar_vista_previa(ruta_archivo, tipo_archivo, ruta_pdf)
                    
                    # Botones de acción (solo para Unidades y Guías)
                    col1, col2 = st.columns(2)
                    with col1:
                        with open(ruta_archivo, "rb") as file:
                            st.download_button(
                                label="Descargar Archivo",
                                data=file,
                                file_name=nombre_archivo,
                                mime=tipo_archivo
                            )
                    with col2:
                        if st.button("Eliminar", key=f"eliminar_{arch_id}"):
                            success, message = eliminar_archivo(arch_id)
                            if success:
                                st.success(message)
                                st.rerun()
                            else:
                                st.error(message)

def alumno_sidebar():
    # Barra lateral simplificada
    user = st.session_state.current_user
    st.sidebar.image(user[6] if user[6] else "https://via.placeholder.com/150", width=150)
    st.sidebar.write(f"**{user[2]} {user[3]}**")
    st.sidebar.write(f"Rol: {user[7]}")
    
    st.sidebar.markdown("---")
    
    if st.sidebar.button("Cerrar Sesión"):
        st.session_state.logged_in = False
        st.session_state.current_user = None
        st.rerun()

def alumno_page():
    user = st.session_state.current_user
    
    # Mostrar barra lateral simplificada
    alumno_sidebar()
    
    # Contenido principal
    st.title(f"Bienvenido, {user[2]} {user[3]}")
    
    # Barra de navegación en el contenido principal
    st.subheader("Navegación")
    col1, col2, col3, col4, col5 = st.columns(5)
    with col1:
        if st.button("Unidades", key="main_unidades"):
            st.session_state.selected_category = "Unidades"
            st.rerun()
    with col2:
        if st.button("Cuestionarios IA", key="main_cuestionarios"):
            st.session_state.selected_category = "Cuestionarios IA"
            st.rerun()
    with col3:
        if st.button("Revisiones", key="main_revisiones"):
            st.session_state.selected_category = "Revisiones"
            st.rerun()
    with col4:
        if st.button("Tutoriales", key="main_tutoriales"):
            st.session_state.selected_category = "Tutoriales"
            st.rerun()
    with col5:
        if st.button("Guías", key="main_guias"):
            st.session_state.selected_category = "Guías"
            st.rerun()
    
    # Contenido según la categoría seleccionada
    category = st.session_state.selected_category
    
    if category == "Cuestionarios IA":
        st.title(f"{category} - En Construcción")
        st.write("Esta sección está en desarrollo. Pronto estará disponible.")
    
    elif category == "Tutoriales":
        st.title("Tutoriales Disponibles")
        
        # Mostrar tutoriales existentes
        tutoriales = obtener_tutoriales()
        
        if not tutoriales:
            st.info("No hay tutoriales disponibles.")
        else:
            for tut in tutoriales:
                tut_id, titulo, descripcion, url_youtube, ruta_video, fecha_creacion, nombre_profesor, apellido_profesor = tut
                
                with st.expander(f"{titulo} - Por: {nombre_profesor} {apellido_profesor}"):
                    st.write(f"**Descripción:** {descripcion}")
                    st.write(f"**Fecha:** {fecha_creacion}")
                    
                    # Mostrar video según el tipo
                    if url_youtube:
                        st.write("**Video de YouTube:**")
                        video_id = extraer_id_youtube(url_youtube)
                        if video_id:
                            st.video(f"https://www.youtube.com/watch?v={video_id}")
                        else:
                            st.error("URL de YouTube no válida")
                        st.link_button("🔗 Ver en YouTube", url_youtube)
                    
                    elif ruta_video and os.path.exists(ruta_video):
                        st.write("**Video:**")
                        st.video(ruta_video)
                        with open(ruta_video, "rb") as file:
                            st.download_button(
                                label="📥 Descargar Video",
                                data=file,
                                file_name=os.path.basename(ruta_video),
                                mime="video/mp4"
                            )
    
    elif category == "Revisiones":
        st.title(f"Archivos Disponibles en {category}")
        
        # Mostrar archivos existentes de la categoría seleccionada
        archivos = obtener_archivos(categoria=category)
        
        if not archivos:
            st.write(f"No hay archivos disponibles en {category}.")
        else:
            for arch in archivos:
                arch_id, nombre_archivo, ruta_archivo, tipo_archivo, fecha_subida, nombre_profesor, apellido_profesor, categoria_archivo, ruta_pdf = arch
                
                # Crear expander para cada archivo
                with st.expander(f"{nombre_archivo} - Subido por: {nombre_profesor} {apellido_profesor}"):
                    st.write(f"Tipo: {tipo_archivo}")
                    st.write(f"Fecha: {fecha_subida}")
                    st.write(f"Categoría: {categoria_archivo}")
                    
                    # Vista previa según tipo de archivo
                    mostrar_vista_previa(ruta_archivo, tipo_archivo, ruta_pdf)
                    
                    # Botón de descarga
                    with open(ruta_archivo, "rb") as file:
                        st.download_button(
                            label="Descargar Archivo",
                            data=file,
                            file_name=nombre_archivo,
                            mime=tipo_archivo
                        )
    else:
        st.title(f"Archivos Disponibles en {category}")
        
        # Mostrar archivos existentes de la categoría seleccionada
        archivos = obtener_archivos(categoria=category)
        
        if not archivos:
            st.write(f"No hay archivos disponibles en {category}.")
        else:
            for arch in archivos:
                arch_id, nombre_archivo, ruta_archivo, tipo_archivo, fecha_subida, nombre_profesor, apellido_profesor, categoria_archivo, ruta_pdf = arch
                
                # Crear expander para cada archivo
                with st.expander(f"{nombre_archivo} - Subido por: {nombre_profesor} {apellido_profesor}"):
                    st.write(f"Tipo: {tipo_archivo}")
                    st.write(f"Fecha: {fecha_subida}")
                    st.write(f"Categoría: {categoria_archivo}")
                    
                    # Vista previa según tipo de archivo
                    mostrar_vista_previa(ruta_archivo, tipo_archivo, ruta_pdf)
                    
                    # Botón de descarga
                    with open(ruta_archivo, "rb") as file:
                        st.download_button(
                            label="Descargar Archivo",
                            data=file,
                            file_name=nombre_archivo,
                            mime=tipo_archivo
                        )

def main_page():
    user = st.session_state.current_user
    rol = user[7]
    
    # Redirigir según el rol
    if rol == 'admin':
        if st.session_state.admin_subpage == "dashboard":
            admin_dashboard()
        elif st.session_state.admin_subpage == "user_management":
            admin_user_management()
    elif rol == 'profesor':
        profesor_page()
    else:  # alumno
        alumno_page()

# Navegación entre páginas
if st.session_state.logged_in:
    main_page()
else:
    if st.session_state.page == "login":
        login_page()
    elif st.session_state.page == "register":
        register_page()
    elif st.session_state.page == "reset_password":
        reset_password_page()
    elif st.session_state.page == "reset_password_confirm":
        reset_password_confirm_page()