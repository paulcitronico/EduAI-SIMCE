import base64
import os
from io import BytesIO
from datetime import datetime

def mostrar_pdf(ruta_archivo):
    with open(ruta_archivo, "rb") as f:
        base64_pdf = base64.b64encode(f.read()).decode('utf-8')
    
    pdf_display = f'<iframe src="data:application/pdf;base64,{base64_pdf}" width="100%" height="800" type="application/pdf"></iframe>'
    return pdf_display

def mostrar_docx_completo(ruta_archivo):
    try:
        from docx import Document
        doc = Document(ruta_archivo)
        
        contenido = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                contenido.append(para.text)
        
        for table in doc.tables:
            contenido.append("---")
            for i, row in enumerate(table.rows):
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text)
                contenido.append(" | ".join(row_data))
            contenido.append("---")
        
        image_count = 0
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_path = os.path.join(os.path.dirname(ruta_archivo), rel.target_ref)
                    if os.path.exists(image_path):
                        contenido.append(f"---\nImagen {image_count + 1}:")
                        contenido.append(image_path)
                        image_count += 1
                except:
                    pass
        
        return contenido
    except Exception as e:
        return [f"No se pudo previsualizar el archivo DOCX: {str(e)}"]

def mostrar_pptx_completo(ruta_archivo):
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        prs = Presentation(ruta_archivo)
        contenido = []
        
        for i, slide in enumerate(prs.slides):
            contenido.append(f"### Diapositiva {i+1}")
            
            slide_title = ""
            for shape in slide.shapes:
                if hasattr(shape, "placeholder_format"):
                    if shape.placeholder_format.type == 1:
                        slide_title = shape.text if shape.text else f"Diapositiva {i+1}"
                        break
            
            if slide_title:
                contenido.append(f"#### {slide_title}")
            
            slide_content = []
            images_in_slide = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_content.append(shape.text)
                
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        
                        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
                        temp_img_path = f"temp_conversion/{timestamp}_slide{i}_img.png"
                        with open(temp_img_path, "wb") as img_file:
                            img_file.write(image_bytes)
                        
                        images_in_slide.append(temp_img_path)
                    except Exception as e:
                        contenido.append(f"Error al procesar imagen: {str(e)}")
                
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for sub_shape in shape.shapes:
                        if hasattr(sub_shape, "text") and sub_shape.text:
                            slide_content.append(sub_shape.text)
                        elif sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                image = sub_shape.image
                                image_bytes = image.blob
                                
                                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
                                temp_img_path = f"temp_conversion/{timestamp}_slide{i}_img.png"
                                with open(temp_img_path, "wb") as img_file:
                                    img_file.write(image_bytes)
                                
                                images_in_slide.append(temp_img_path)
                            except Exception as e:
                                contenido.append(f"Error al procesar imagen en grupo: {str(e)}")
            
            if slide_content:
                for text in slide_content:
                    if text.strip():
                        contenido.append(text)
            
            if images_in_slide:
                contenido.append("---\nImágenes en la diapositiva:")
                contenido.extend(images_in_slide)
            
            contenido.append("---")
        
        return contenido
    except Exception as e:
        return [f"No se pudo previsualizar el archivo PPTX: {str(e)}"]

def mostrar_imagen(ruta_archivo):
    try:
        from PIL import Image
        image = Image.open(ruta_archivo)
        return image
    except Exception as e:
        return f"No se pudo previsualizar la imagen: {str(e)}"

def mostrar_vista_previa(ruta_archivo, tipo_archivo, ruta_pdf=None):
    import streamlit as st
    
    if tipo_archivo == 'application/pdf':
        st.write("Vista previa del PDF:")
        st.markdown(mostrar_pdf(ruta_archivo), unsafe_allow_html=True)
    elif tipo_archivo in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 
                         'application/msword',
                         'application/vnd.ms-word.document.macroEnabled.12']:
        st.write("Vista previa del documento:")
        contenido = mostrar_docx_completo(ruta_archivo)
        for item in contenido:
            if item.startswith("temp_conversion/") and os.path.exists(item):
                st.image(item, width=300)
            else:
                st.write(item)
    elif tipo_archivo in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         'application/vnd.ms-powerpoint',
                         'application/vnd.ms-powerpoint.presentation.macroEnabled.12']:
        if ruta_pdf and os.path.exists(ruta_pdf):
            st.write("Vista previa de la presentación (convertida a PDF):")
            st.markdown(mostrar_pdf(ruta_pdf), unsafe_allow_html=True)
        else:
            st.write("Vista previa de la presentación:")
            contenido = mostrar_pptx_completo(ruta_archivo)
            for item in contenido:
                if item.startswith("temp_conversion/") and os.path.exists(item):
                    st.image(item, width=300)
                else:
                    st.write(item)
    elif tipo_archivo.startswith('image/'):
        st.write("Vista previa de la imagen:")
        imagen = mostrar_imagen(ruta_archivo)
        if isinstance(imagen, str):
            st.error(imagen)
        else:
            st.image(imagen, width=800)
    else:
        st.write("Vista previa no disponible para este tipo de archivo")